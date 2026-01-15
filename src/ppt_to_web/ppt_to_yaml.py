import subprocess
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE


def _extract_text_from_shape(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(
        [run.text for para in shape.text_frame.paragraphs for run in para.runs]
    )


def _is_highlighted(shape) -> bool:
    return hasattr(shape, "fill") and shape.fill.type != 0


def _map_chart_type(chart_type: XL_CHART_TYPE) -> str:
    """Map python-pptx XL_CHART_TYPE to ECharts chart type string."""
    chart_type_mapping = {
        # Bar charts (horizontal bars in PowerPoint)
        XL_CHART_TYPE.BAR_CLUSTERED: "bar",
        XL_CHART_TYPE.BAR_STACKED: "bar",
        XL_CHART_TYPE.BAR_STACKED_100: "bar",
        # Column charts (vertical bars - most common)
        XL_CHART_TYPE.COLUMN_CLUSTERED: "bar",
        XL_CHART_TYPE.COLUMN_STACKED: "bar",
        XL_CHART_TYPE.COLUMN_STACKED_100: "bar",
        # Line charts
        XL_CHART_TYPE.LINE: "line",
        XL_CHART_TYPE.LINE_MARKERS: "line",
        XL_CHART_TYPE.LINE_STACKED: "line",
        XL_CHART_TYPE.LINE_STACKED_100: "line",
        XL_CHART_TYPE.LINE_MARKERS_STACKED: "line",
        XL_CHART_TYPE.LINE_MARKERS_STACKED_100: "line",
        # Pie charts
        XL_CHART_TYPE.PIE: "pie",
        XL_CHART_TYPE.PIE_EXPLODED: "pie",
        XL_CHART_TYPE.DOUGHNUT: "pie",
        XL_CHART_TYPE.DOUGHNUT_EXPLODED: "pie",
        # Area charts (ECharts line with areaStyle)
        XL_CHART_TYPE.AREA: "line",
        XL_CHART_TYPE.AREA_STACKED: "line",
        XL_CHART_TYPE.AREA_STACKED_100: "line",
        # Scatter/XY charts
        XL_CHART_TYPE.XY_SCATTER: "scatter",
        XL_CHART_TYPE.XY_SCATTER_LINES: "scatter",
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS: "scatter",
        XL_CHART_TYPE.XY_SCATTER_SMOOTH: "scatter",
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS: "scatter",
        # Radar charts
        XL_CHART_TYPE.RADAR: "radar",
        XL_CHART_TYPE.RADAR_FILLED: "radar",
        XL_CHART_TYPE.RADAR_MARKERS: "radar",
    }
    return chart_type_mapping.get(chart_type, "bar")


def _extract_chart_title(chart) -> str:
    """Extract chart title safely."""
    if not chart.has_title:
        return ""
    try:
        return chart.chart_title.text_frame.text
    except Exception:
        return ""


def _extract_chart_categories(chart) -> list[str]:
    """Extract chart categories (x-axis labels)."""
    try:
        if chart.plots and chart.plots[0].categories:
            return [str(cat) for cat in chart.plots[0].categories]
    except Exception as e:
        print(f"Warning: Could not extract categories: {e}")
    return []


def _extract_chart_series(chart) -> list[dict]:
    """Extract series data from chart."""
    series_data = []
    for idx, series in enumerate(chart.series):
        series_name = str(series.name) if series.name else f"Series {idx + 1}"
        series_values = []
        try:
            if series.values:
                series_values = [float(v) if v is not None else 0 for v in series.values]
        except Exception:
            pass
        series_data.append({"name": series_name, "data": series_values})
    return series_data


HORIZONTAL_BAR_TYPES = {
    XL_CHART_TYPE.BAR_CLUSTERED,
    XL_CHART_TYPE.BAR_STACKED,
    XL_CHART_TYPE.BAR_STACKED_100,
}


def _extract_chart(shape, slide_idx: int, shape_idx: int) -> dict | None:
    """Extract chart data from a shape containing a chart."""
    if not shape.has_chart:
        return None

    try:
        chart = shape.chart
        chart_type_enum = chart.chart_type

        series_data = _extract_chart_series(chart)
        if not series_data or all(not s["data"] for s in series_data):
            return None

        return {
            "type": "chart",
            "chart_type": _map_chart_type(chart_type_enum),
            "title": _extract_chart_title(chart),
            "categories": _extract_chart_categories(chart),
            "series": series_data,
            "is_stacked": "STACKED" in chart_type_enum.name,
            "is_horizontal": chart_type_enum in HORIZONTAL_BAR_TYPES,
            "is_area": "AREA" in chart_type_enum.name,
            "chart_id": f"chart_{slide_idx}_{shape_idx}",
        }

    except Exception as e:
        print(f"Warning: Failed to extract chart from slide {slide_idx}, shape {shape_idx}: {e}")
        return None


def _convert_with_libreoffice(input_path: Path, output_dir: Path) -> Path | None:
    """Convert image using LibreOffice (best for WMF/EMF)."""
    try:
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", str(output_dir),
                str(input_path),
            ],
            capture_output=True,
            timeout=60,
        )
        if result.returncode == 0:
            output_path = output_dir / f"{input_path.stem}.png"
            if output_path.exists():
                return output_path
    except Exception:
        pass
    return None


def _get_image_dimensions(filepath: Path) -> tuple[int, int, float]:
    """Get image dimensions using wand or return defaults."""
    width, height, aspect_ratio = 0, 0, 1.0
    try:
        from wand.image import Image as WandImage
        with WandImage(filename=str(filepath)) as img:
            width = img.width
            height = img.height
            if height > 0:
                aspect_ratio = width / height
    except Exception:
        pass
    return width, height, aspect_ratio


WEB_IMAGE_FORMATS = ("png", "jpg", "jpeg", "gif", "webp")
VECTOR_IMAGE_FORMATS = ("wmf", "emf", "wmz", "emz")


def _make_media_result(path: str, width: int = 0, height: int = 0) -> dict:
    """Create a standardized media result dictionary."""
    aspect_ratio = width / height if height > 0 else 1.0
    return {"path": path, "width": width, "height": height, "aspect_ratio": aspect_ratio}


def _process_web_image(image_bytes: bytes, output_path: Path) -> dict | None:
    """Process web-compatible image formats using Wand."""
    try:
        from wand.image import Image as WandImage
        from wand.color import Color

        with WandImage(blob=image_bytes) as img:
            img.trim(color=Color('white'), fuzz=0)
            img.trim(fuzz=0)
            img.format = 'png'
            img.save(filename=str(output_path))
            return _make_media_result(f"media/{output_path.name}", img.width, img.height)
    except Exception:
        return None


def _process_vector_image(
    image_bytes: bytes, ext: str, slide_idx: int, shape_idx: int, media_dir: Path
) -> dict:
    """Process vector image formats (WMF/EMF) using LibreOffice."""
    temp_filename = f"slide_{slide_idx}_shape_{shape_idx}.{ext}"
    temp_filepath = media_dir / temp_filename
    png_filename = f"slide_{slide_idx}_shape_{shape_idx}.png"

    with open(temp_filepath, "wb") as f:
        f.write(image_bytes)

    converted = _convert_with_libreoffice(temp_filepath, media_dir)
    if converted and converted.exists():
        temp_filepath.unlink(missing_ok=True)
        width, height, aspect_ratio = _get_image_dimensions(converted)
        return _make_media_result(f"media/{png_filename}", width, height)

    print(f"Warning: LibreOffice conversion failed for slide {slide_idx}, shape {shape_idx}")
    return _make_media_result(f"media/{temp_filename}")


def _extract_media(
    shape, slide_idx: int, shape_idx: int, media_dir: Path
) -> dict | None:
    """Extract and process media from a shape."""
    if not hasattr(shape, "image"):
        return None

    try:
        image_bytes = shape.image.blob
        ext = shape.image.ext.lower()
        png_filename = f"slide_{slide_idx}_shape_{shape_idx}.png"
        png_filepath = media_dir / png_filename

        # Try Wand for web-compatible formats
        if ext in WEB_IMAGE_FORMATS:
            result = _process_web_image(image_bytes, png_filepath)
            if result:
                return result
            print(f"Warning: Wand processing failed for slide {slide_idx}, shape {shape_idx}")

        # Try LibreOffice for vector formats or if Wand failed
        if ext in VECTOR_IMAGE_FORMATS or not png_filepath.exists():
            return _process_vector_image(image_bytes, ext, slide_idx, shape_idx, media_dir)

        # Fallback: save original format
        filename = f"slide_{slide_idx}_shape_{shape_idx}.{ext}"
        with open(media_dir / filename, "wb") as f:
            f.write(image_bytes)
        return _make_media_result(f"media/{filename}")

    except Exception as e:
        print(f"Warning: Failed to extract media from slide {slide_idx}, shape {shape_idx}: {e}")
        return None


def ppt_to_yaml(pptx_path: str, yaml_output_dir: str) -> str:
    pptx_file = Path(pptx_path)
    yaml_dir = Path(yaml_output_dir)
    yaml_dir.mkdir(parents=True, exist_ok=True)

    media_dir = yaml_dir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_file))

    slides_data = []
    highlighted_sections = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": slide_idx + 1,
            "title": "",
            "content": [],
            "media": [],
            "is_highlighted": False,
            "layout": slide.slide_layout.name if slide.slide_layout else "",
        }

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = _extract_text_from_shape(shape).strip()
                if text:
                    if shape_idx == 0 and not slide_data["title"]:
                        slide_data["title"] = text
                    else:
                        slide_data["content"].append({"type": "text", "value": text})

            # Check for chart BEFORE image (charts may also have image representations)
            if shape.has_chart:
                chart_data = _extract_chart(shape, slide_idx, shape_idx)
                if chart_data:
                    slide_data["media"].append(chart_data)
            elif hasattr(shape, "image"):
                media_info = _extract_media(shape, slide_idx, shape_idx, media_dir)
                if media_info:
                    media_item = {"type": "image"}
                    media_item.update(media_info)
                    slide_data["media"].append(media_item)

            if _is_highlighted(shape):
                slide_data["is_highlighted"] = True

        if slide_data["title"] or slide_data["content"] or slide_data["media"]:
            slides_data.append(slide_data)
            if slide_data["is_highlighted"]:
                highlighted_sections.append(
                    {
                        "slide_number": slide_idx + 1,
                        "title": slide_data["title"],
                        "content": slide_data["content"][:3],
                    }
                )

    # Extract cover title from first slide
    cover_title = pptx_file.stem
    if slides_data and slides_data[0].get("title"):
        cover_title = slides_data[0]["title"].replace("\n", " ").strip()

    output_data = {
        "title": pptx_file.stem,
        "cover_title": cover_title,
        "hero_image": None,
        "slides": slides_data,
        "highlighted_sections": highlighted_sections,
        "total_slides": len(slides_data),
    }

    yaml_path = yaml_dir / f"{pptx_file.stem}.yaml"
    with open(yaml_path, "w", encoding="utf-8") as f:
        yaml.dump(
            output_data,
            f,
            allow_unicode=True,
            default_flow_style=False,
            sort_keys=False,
        )

    return str(yaml_path)
