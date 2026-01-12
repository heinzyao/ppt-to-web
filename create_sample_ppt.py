#!/usr/bin/env python
"""
Create a sample PowerPoint presentation with images for testing.

Usage:
    uv run python create_sample_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


def create_sample_ppt():
    prs = Presentation()

    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle

        return slide

    def add_content_slide(title, content_items, is_highlighted=False):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if is_highlighted:
            try:
                slide.shapes.title.fill.solid()
                slide.shapes.title.fill.fore_color.rgb = RGBColor(114, 47, 55)
            except:
                pass

        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

        return slide

    def add_image_slide(title, image_path, caption=""):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1.5)

        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8))

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title

        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1), Inches(5), Inches(8), Inches(0.5)
            )
            tf = caption_box.text_frame
            tf.text = caption

        return slide

    def add_mixed_slide(title, image_path, content_items):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title

        pic = slide.shapes.add_picture(
            image_path, Inches(1), Inches(1.5), width=Inches(4)
        )

        content_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(1.5), Inches(3.5), Inches(3)
        )
        tf = content_box.text_frame
        tf.clear()

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

        return slide

    add_title_slide(
        "Strategic Digital Transformation", "A Path to Organizational Excellence"
    )

    add_content_slide(
        "Executive Summary",
        [
            "Digital transformation is no longer optional",
            "Companies must adapt to changing market dynamics",
            "This presentation outlines key strategic priorities",
        ],
        is_highlighted=True,
    )

    add_content_slide(
        "Market Analysis",
        [
            "75% of companies report accelerated digital adoption",
            "Cloud computing market growth: 18% YoY",
            "AI integration projected to add $15.7T to global GDP by 2030",
        ],
        is_highlighted=True,
    )

    add_mixed_slide(
        "Revenue Growth Performance",
        "sample_images/growth_chart.png",
        [
            "Strong year-over-year growth",
            "45% increase in 2021",
            "95% projected by 2024",
            "Outperformed industry average",
        ],
    )

    add_image_slide(
        "Market Share Distribution",
        "sample_images/market_share.png",
        "Source: Internal Market Analysis Report 2024",
    )

    add_mixed_slide(
        "Key Success Indicators",
        "sample_images/target_icon.png",
        [
            "Customer Satisfaction",
            "Net Promoter Score > 50",
            "Customer Retention > 90%",
            "Response Time < 24hrs",
        ],
    )

    add_content_slide(
        "Implementation Roadmap",
        [
            "Phase 1: Foundation (Q1-Q2)",
            "   • Infrastructure modernization",
            "   • Team capability building",
            "Phase 2: Acceleration (Q3-Q4)",
            "   • Core system integration",
            "   • Pilot program launch",
        ],
    )

    add_image_slide(
        "Strategic Roadmap Visualization",
        "sample_images/roadmap_diagram.png",
        "Four-phase approach to digital transformation",
    )

    add_content_slide(
        "Success Metrics",
        [
            "Customer satisfaction: +15%",
            "Cost reduction: 20%",
            "Time-to-market: 30% faster",
            "Revenue growth: 12% YoY",
        ],
        is_highlighted=True,
    )

    add_content_slide(
        "Conclusion",
        [
            "Strategic digital transformation drives sustainable growth",
            "Leadership commitment is critical for success",
            "Continuous adaptation to market changes is essential",
        ],
    )

    output_path = "sample_presentation.pptx"
    prs.save(output_path)
    print(f"✓ Sample presentation created: {output_path}")
    print(f"  • 10 slides including images and mixed content")
    print(f"  • 3 highlighted sections for cover focus")
    return output_path


if __name__ == "__main__":
    create_sample_ppt()
