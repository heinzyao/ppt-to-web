"""Tests for ppt_to_yaml module."""

from unittest.mock import MagicMock, PropertyMock, patch

import pytest
from pptx.enum.chart import XL_CHART_TYPE

from ppt_to_web.ppt_to_yaml import (
    _extract_chart,
    _extract_chart_categories,
    _extract_chart_series,
    _extract_chart_title,
    _extract_text_from_shape,
    _is_highlighted,
    _make_media_result,
    _map_chart_type,
    ppt_to_yaml,
)


# --- _extract_text_from_shape ---


class TestExtractTextFromShape:
    def test_no_text_frame(self):
        shape = MagicMock()
        shape.has_text_frame = False
        assert _extract_text_from_shape(shape) == ""

    def test_single_paragraph_single_run(self):
        run = MagicMock()
        run.text = "Hello World"
        para = MagicMock()
        para.runs = [run]
        shape = MagicMock()
        shape.has_text_frame = True
        shape.text_frame.paragraphs = [para]
        assert _extract_text_from_shape(shape) == "Hello World"

    def test_multiple_paragraphs(self):
        run1 = MagicMock(text="Line 1")
        run2 = MagicMock(text="Line 2")
        para1 = MagicMock(runs=[run1])
        para2 = MagicMock(runs=[run2])
        shape = MagicMock()
        shape.has_text_frame = True
        shape.text_frame.paragraphs = [para1, para2]
        assert _extract_text_from_shape(shape) == "Line 1\nLine 2"

    def test_multiple_runs_in_paragraph(self):
        """Each run becomes a separate line due to join with newline."""
        run1 = MagicMock(text="Hello ")
        run2 = MagicMock(text="World")
        para = MagicMock(runs=[run1, run2])
        shape = MagicMock()
        shape.has_text_frame = True
        shape.text_frame.paragraphs = [para]
        # Implementation joins all runs across all paragraphs with \n
        assert _extract_text_from_shape(shape) == "Hello \nWorld"

    def test_empty_paragraphs(self):
        para = MagicMock(runs=[])
        shape = MagicMock()
        shape.has_text_frame = True
        shape.text_frame.paragraphs = [para]
        assert _extract_text_from_shape(shape) == ""


# --- _is_highlighted ---


class TestIsHighlighted:
    def test_no_fill_attribute(self):
        shape = MagicMock(spec=[])  # no fill attribute
        assert _is_highlighted(shape) is False

    def test_fill_type_zero(self):
        shape = MagicMock()
        shape.fill.type = 0
        assert _is_highlighted(shape) is False

    def test_fill_type_nonzero(self):
        shape = MagicMock()
        shape.fill.type = 1
        assert _is_highlighted(shape) is True


# --- _map_chart_type ---


class TestMapChartType:
    @pytest.mark.parametrize(
        "pptx_type,expected",
        [
            (XL_CHART_TYPE.COLUMN_CLUSTERED, "bar"),
            (XL_CHART_TYPE.COLUMN_STACKED, "bar"),
            (XL_CHART_TYPE.BAR_CLUSTERED, "bar"),
            (XL_CHART_TYPE.LINE, "line"),
            (XL_CHART_TYPE.LINE_MARKERS, "line"),
            (XL_CHART_TYPE.PIE, "pie"),
            (XL_CHART_TYPE.DOUGHNUT, "pie"),
            (XL_CHART_TYPE.AREA, "line"),
            (XL_CHART_TYPE.XY_SCATTER, "scatter"),
            (XL_CHART_TYPE.RADAR, "radar"),
            (XL_CHART_TYPE.RADAR_FILLED, "radar"),
        ],
    )
    def test_known_types(self, pptx_type, expected):
        assert _map_chart_type(pptx_type) == expected

    def test_unknown_type_defaults_to_bar(self):
        assert _map_chart_type("UNKNOWN_TYPE") == "bar"


# --- _make_media_result ---


class TestMakeMediaResult:
    def test_basic(self):
        result = _make_media_result("media/img.png", 800, 600)
        assert result["path"] == "media/img.png"
        assert result["width"] == 800
        assert result["height"] == 600
        assert result["aspect_ratio"] == pytest.approx(800 / 600)

    def test_zero_height(self):
        result = _make_media_result("media/img.png", 800, 0)
        assert result["aspect_ratio"] == 1.0

    def test_defaults(self):
        result = _make_media_result("media/img.png")
        assert result["width"] == 0
        assert result["height"] == 0
        assert result["aspect_ratio"] == 1.0


# --- _extract_chart_title ---


class TestExtractChartTitle:
    def test_no_title(self):
        chart = MagicMock()
        chart.has_title = False
        assert _extract_chart_title(chart) == ""

    def test_with_title(self):
        chart = MagicMock()
        chart.has_title = True
        chart.chart_title.text_frame.text = "Sales Report"
        assert _extract_chart_title(chart) == "Sales Report"

    def test_title_extraction_error(self):
        chart = MagicMock()
        chart.has_title = True
        type(chart).chart_title = PropertyMock(side_effect=Exception("broken"))
        assert _extract_chart_title(chart) == ""


# --- _extract_chart_categories ---


class TestExtractChartCategories:
    def test_with_categories(self):
        chart = MagicMock()
        chart.plots.__getitem__(0).categories = ["Q1", "Q2", "Q3"]
        chart.plots.__bool__ = lambda self: True
        assert _extract_chart_categories(chart) == ["Q1", "Q2", "Q3"]

    def test_no_plots(self):
        chart = MagicMock()
        chart.plots = []
        assert _extract_chart_categories(chart) == []

    def test_exception_returns_empty(self):
        chart = MagicMock()
        type(chart).plots = PropertyMock(side_effect=Exception("broken"))
        assert _extract_chart_categories(chart) == []


# --- _extract_chart_series ---


class TestExtractChartSeries:
    def test_single_series(self):
        series = MagicMock()
        series.name = "Revenue"
        series.values = [10.0, 20.0, 30.0]
        chart = MagicMock()
        chart.series = [series]
        result = _extract_chart_series(chart)
        assert len(result) == 1
        assert result[0]["name"] == "Revenue"
        assert result[0]["data"] == [10.0, 20.0, 30.0]

    def test_multiple_series(self):
        s1 = MagicMock(name="A", values=[1.0, 2.0])
        s1.name = "A"
        s2 = MagicMock(name="B", values=[3.0, 4.0])
        s2.name = "B"
        chart = MagicMock()
        chart.series = [s1, s2]
        result = _extract_chart_series(chart)
        assert len(result) == 2

    def test_none_values_become_zero(self):
        series = MagicMock()
        series.name = "Data"
        series.values = [1.0, None, 3.0]
        chart = MagicMock()
        chart.series = [series]
        result = _extract_chart_series(chart)
        assert result[0]["data"] == [1.0, 0.0, 3.0]

    def test_unnamed_series(self):
        series = MagicMock()
        series.name = None
        series.values = [1.0]
        chart = MagicMock()
        chart.series = [series]
        result = _extract_chart_series(chart)
        assert result[0]["name"] == "Series 1"

    def test_empty_series(self):
        chart = MagicMock()
        chart.series = []
        assert _extract_chart_series(chart) == []


# --- _extract_chart ---


class TestExtractChart:
    def test_no_chart(self):
        shape = MagicMock()
        shape.has_chart = False
        assert _extract_chart(shape, 0, 0) is None

    def test_valid_chart(self):
        series = MagicMock()
        series.name = "Sales"
        series.values = [10.0, 20.0]
        chart = MagicMock()
        chart.chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart.has_title = True
        chart.chart_title.text_frame.text = "Revenue"
        chart.plots.__getitem__(0).categories = ["Q1", "Q2"]
        chart.plots.__bool__ = lambda self: True
        chart.series = [series]

        shape = MagicMock()
        shape.has_chart = True
        shape.chart = chart

        result = _extract_chart(shape, 1, 2)
        assert result is not None
        assert result["type"] == "chart"
        assert result["chart_type"] == "bar"
        assert result["title"] == "Revenue"
        assert result["chart_id"] == "chart_1_2"
        assert result["is_stacked"] is False
        assert result["is_horizontal"] is False

    def test_stacked_bar(self):
        series = MagicMock()
        series.name = "Data"
        series.values = [1.0]
        chart = MagicMock()
        chart.chart_type = XL_CHART_TYPE.BAR_STACKED
        chart.has_title = False
        chart.plots.__getitem__(0).categories = ["A"]
        chart.plots.__bool__ = lambda self: True
        chart.series = [series]

        shape = MagicMock()
        shape.has_chart = True
        shape.chart = chart

        result = _extract_chart(shape, 0, 0)
        assert result["is_stacked"] is True
        assert result["is_horizontal"] is True

    def test_area_chart(self):
        series = MagicMock()
        series.name = "Data"
        series.values = [1.0]
        chart = MagicMock()
        chart.chart_type = XL_CHART_TYPE.AREA
        chart.has_title = False
        chart.plots.__getitem__(0).categories = []
        chart.plots.__bool__ = lambda self: True
        chart.series = [series]

        shape = MagicMock()
        shape.has_chart = True
        shape.chart = chart

        result = _extract_chart(shape, 0, 0)
        assert result["chart_type"] == "line"
        assert result["is_area"] is True

    def test_empty_series_returns_none(self):
        chart = MagicMock()
        chart.chart_type = XL_CHART_TYPE.PIE
        chart.series = []

        shape = MagicMock()
        shape.has_chart = True
        shape.chart = chart

        assert _extract_chart(shape, 0, 0) is None

    def test_exception_returns_none(self):
        shape = MagicMock()
        shape.has_chart = True
        type(shape).chart = PropertyMock(side_effect=Exception("broken"))
        assert _extract_chart(shape, 0, 0) is None


# --- ppt_to_yaml (integration) ---


class TestPptToYaml:
    def _make_presentation(self, slides_config):
        """Helper to build a mock Presentation.

        slides_config: list of dicts, each with:
          - shapes: list of shape dicts with keys like 'text', 'has_chart', 'has_image'
        """
        prs = MagicMock()
        slides = []
        for slide_cfg in slides_config:
            slide = MagicMock()
            slide.slide_layout.name = "Title Slide"
            shapes = []
            for shape_cfg in slide_cfg.get("shapes", []):
                shape = MagicMock()
                # Text
                if "text" in shape_cfg:
                    shape.has_text_frame = True
                    run = MagicMock(text=shape_cfg["text"])
                    para = MagicMock(runs=[run])
                    shape.text_frame.paragraphs = [para]
                else:
                    shape.has_text_frame = False

                # Chart
                shape.has_chart = shape_cfg.get("has_chart", False)

                # Image - remove 'image' attribute if not an image shape
                if not shape_cfg.get("has_image", False):
                    del shape.image

                # Highlight
                shape.fill.type = 1 if shape_cfg.get("highlighted", False) else 0

                shapes.append(shape)
            slide.shapes = shapes
            slides.append(slide)
        prs.slides = slides
        return prs

    @patch("ppt_to_web.ppt_to_yaml.Presentation")
    def test_basic_text_extraction(self, mock_prs_cls, tmp_path):
        mock_prs_cls.return_value = self._make_presentation(
            [
                {"shapes": [{"text": "Slide Title"}, {"text": "Body content"}]},
            ]
        )

        pptx_path = tmp_path / "test.pptx"
        pptx_path.touch()
        output_dir = tmp_path / "output"

        yaml_path = ppt_to_yaml(str(pptx_path), str(output_dir))

        import yaml

        with open(yaml_path) as f:
            data = yaml.safe_load(f)

        assert data["title"] == "test"
        assert data["cover_title"] == "Slide Title"
        assert data["total_slides"] == 1
        assert len(data["slides"]) == 1
        assert data["slides"][0]["title"] == "Slide Title"
        assert data["slides"][0]["content"][0]["value"] == "Body content"

    @patch("ppt_to_web.ppt_to_yaml.Presentation")
    def test_empty_presentation(self, mock_prs_cls, tmp_path):
        mock_prs_cls.return_value = self._make_presentation([{"shapes": []}])

        pptx_path = tmp_path / "empty.pptx"
        pptx_path.touch()
        output_dir = tmp_path / "output"

        yaml_path = ppt_to_yaml(str(pptx_path), str(output_dir))

        import yaml

        with open(yaml_path) as f:
            data = yaml.safe_load(f)

        assert data["total_slides"] == 0
        assert data["slides"] == []
        assert data["cover_title"] == "empty"

    @patch("ppt_to_web.ppt_to_yaml.Presentation")
    def test_highlighted_slide(self, mock_prs_cls, tmp_path):
        mock_prs_cls.return_value = self._make_presentation(
            [
                {
                    "shapes": [
                        {"text": "Important", "highlighted": True},
                    ]
                },
            ]
        )

        pptx_path = tmp_path / "highlight.pptx"
        pptx_path.touch()
        output_dir = tmp_path / "output"

        yaml_path = ppt_to_yaml(str(pptx_path), str(output_dir))

        import yaml

        with open(yaml_path) as f:
            data = yaml.safe_load(f)

        assert data["slides"][0]["is_highlighted"] is True
        assert len(data["highlighted_sections"]) == 1

    @patch("ppt_to_web.ppt_to_yaml.Presentation")
    def test_multiple_slides(self, mock_prs_cls, tmp_path):
        mock_prs_cls.return_value = self._make_presentation(
            [
                {"shapes": [{"text": "Title 1"}, {"text": "Content 1"}]},
                {"shapes": [{"text": "Title 2"}, {"text": "Content 2"}]},
            ]
        )

        pptx_path = tmp_path / "multi.pptx"
        pptx_path.touch()
        output_dir = tmp_path / "output"

        yaml_path = ppt_to_yaml(str(pptx_path), str(output_dir))

        import yaml

        with open(yaml_path) as f:
            data = yaml.safe_load(f)

        assert data["total_slides"] == 2
        assert data["slides"][0]["title"] == "Title 1"
        assert data["slides"][1]["title"] == "Title 2"

    @patch("ppt_to_web.ppt_to_yaml.Presentation")
    def test_creates_media_directory(self, mock_prs_cls, tmp_path):
        mock_prs_cls.return_value = self._make_presentation(
            [{"shapes": [{"text": "Test"}]}]
        )

        pptx_path = tmp_path / "test.pptx"
        pptx_path.touch()
        output_dir = tmp_path / "output"

        ppt_to_yaml(str(pptx_path), str(output_dir))

        assert (output_dir / "media").is_dir()

    @patch("ppt_to_web.ppt_to_yaml.Presentation")
    def test_output_yaml_structure(self, mock_prs_cls, tmp_path):
        mock_prs_cls.return_value = self._make_presentation(
            [{"shapes": [{"text": "Title"}]}]
        )

        pptx_path = tmp_path / "struct.pptx"
        pptx_path.touch()
        output_dir = tmp_path / "output"

        yaml_path = ppt_to_yaml(str(pptx_path), str(output_dir))

        import yaml

        with open(yaml_path) as f:
            data = yaml.safe_load(f)

        expected_keys = {"title", "cover_title", "hero_image", "slides", "highlighted_sections", "total_slides"}
        assert set(data.keys()) == expected_keys
        assert data["hero_image"] is None
